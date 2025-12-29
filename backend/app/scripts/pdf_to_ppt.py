#!/usr/bin/env python3
"""
PDF 到 PowerPoint 转换脚本
图像级转换：将 PDF 每页渲染为高质量图片，然后插入到 PPT 中
优势：
- 完美保留 PDF 视觉效果
- 不依赖 LibreOffice
- 实现简单、稳定可靠
- 支持高 DPI 渲染（可配置）
"""
import argparse
import os
import sys
import time
import traceback
import tempfile
import shutil
from typing import List, Tuple
from pathlib import Path

try:
    from pdf2image import convert_from_path
    from PIL import Image
except ImportError:
    print("[ERROR] 缺少必要的依赖库")
    print("[ERROR] 请安装: pip install pdf2image Pillow")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("[ERROR] 缺少必要的依赖库")
    print("[ERROR] 请安装: pip install python-pptx")
    sys.exit(1)


def _check_poppler() -> bool:
    """
    检查 poppler 是否可用（pdf2image 的依赖）
    Windows 需要单独安装 poppler
    """
    try:
        # 尝试转换一个测试 PDF（实际上我们只是检查命令是否存在）
        import subprocess
        result = subprocess.run(
            ["pdftoppm", "-h"],
            capture_output=True,
            timeout=5
        )
        return True
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False


def pdf_to_images(pdf_path: str, dpi: int = 200) -> List[Image.Image]:
    """
    将 PDF 的每一页转换为图片
    
    参数:
        pdf_path: PDF 文件路径
        dpi: 图像分辨率，默认 200（较高质量），可选 150, 300
    
    返回:
        图片列表（PIL Image 对象）
    """
    print(f"[INFO] 正在将 PDF 转换为图片 (DPI={dpi})...")
    
    try:
        # Windows 平台特殊处理
        if sys.platform == "win32":
            # 尝试查找 poppler
            poppler_path = None
            possible_paths = [
                r"C:\Program Files\poppler\Library\bin",
                r"C:\Program Files (x86)\poppler\Library\bin",
                os.path.join(os.path.dirname(sys.executable), "Library", "bin"),
            ]
            
            for path in possible_paths:
                if os.path.exists(path) and os.path.exists(os.path.join(path, "pdftoppm.exe")):
                    poppler_path = path
                    print(f"[INFO] 找到 poppler: {poppler_path}")
                    break
            
            if poppler_path:
                images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
            else:
                print("[WARN] 未找到 poppler，尝试使用系统 PATH")
                images = convert_from_path(pdf_path, dpi=dpi)
        else:
            images = convert_from_path(pdf_path, dpi=dpi)
        
        print(f"[SUCCESS] 成功转换 {len(images)} 页")
        return images
        
    except Exception as e:
        print(f"[ERROR] PDF 转图片失败: {str(e)}")
        print("[ERROR] 请确保已安装 poppler:")
        print("[ERROR]   - Windows: 下载 https://github.com/oschwartz10612/poppler-windows/releases/")
        print("[ERROR]   - Linux: sudo apt-get install poppler-utils")
        print("[ERROR]   - macOS: brew install poppler")
        raise


def create_ppt_from_images(images: List[Image.Image], output_path: str, 
                          slide_width: float = 10, slide_height: float = 7.5) -> bool:
    """
    从图片列表创建 PowerPoint 文件
    
    参数:
        images: PIL Image 对象列表
        output_path: 输出 PPTX 文件路径
        slide_width: 幻灯片宽度（英寸），默认 10
        slide_height: 幻灯片高度（英寸），默认 7.5
    
    返回:
        是否成功
    """
    try:
        print(f"[INFO] 正在创建 PowerPoint 文件...")
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)
        
        # 为每张图片创建一个幻灯片
        for i, image in enumerate(images, 1):
            print(f"[INFO] 正在添加第 {i}/{len(images)} 页...")
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # 6 = 空白布局
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 保存临时图片文件
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                temp_image_path = tmp_file.name
                image.save(temp_image_path, 'PNG')
            
            try:
                # 计算图片缩放比例以适应幻灯片
                img_width, img_height = image.size
                slide_width_px = prs.slide_width
                slide_height_px = prs.slide_height
                
                # 计算缩放比例（保持宽高比）
                width_ratio = slide_width_px / img_width
                height_ratio = slide_height_px / img_height
                scale_ratio = min(width_ratio, height_ratio)
                
                # 计算最终尺寸和位置（居中）
                final_width = int(img_width * scale_ratio)
                final_height = int(img_height * scale_ratio)
                left = (slide_width_px - final_width) // 2
                top = (slide_height_px - final_height) // 2
                
                # 添加图片到幻灯片
                slide.shapes.add_picture(
                    temp_image_path,
                    left, top,
                    width=final_width,
                    height=final_height
                )
                
            finally:
                # 清理临时文件
                try:
                    os.unlink(temp_image_path)
                except:
                    pass
        
        # 保存 PowerPoint 文件
        print(f"[INFO] 正在保存 PowerPoint 文件...")
        prs.save(output_path)
        
        print(f"[SUCCESS] PowerPoint 文件已创建: {output_path}")
        return True
        
    except Exception as e:
        print(f"[ERROR] 创建 PowerPoint 失败: {str(e)}")
        traceback.print_exc()
        return False


def pdf_to_ppt(pdf_path: str, ppt_path: str, dpi: int = 200) -> bool:
    """
    将 PDF 转换为 PowerPoint（图像级转换）
    
    参数:
        pdf_path: 输入 PDF 文件路径
        ppt_path: 输出 PPTX 文件路径
        dpi: 图像分辨率（150=标准, 200=高质量, 300=超高质量）
    
    返回:
        是否成功
    """
    try:
        print(f"[INFO] ========== PDF 转 PPT 开始 ==========")
        print(f"[INFO] 转换方案: 图像级转换（每页渲染为图片）")
        start_time = time.time()
        
        # 验证输入文件
        pdf_path = os.path.abspath(pdf_path)
        if not os.path.exists(pdf_path):
            print(f"[ERROR] 输入文件不存在: {pdf_path}")
            return False
        
        print(f"[INFO] 输入PDF: {pdf_path}")
        print(f"[INFO] PDF文件大小: {os.path.getsize(pdf_path) / (1024*1024):.2f} MB")
        print(f"[INFO] 图像质量: {dpi} DPI")
        
        # 准备输出目录
        ppt_path = os.path.abspath(ppt_path)
        output_dir = os.path.dirname(ppt_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        
        # 第一步：PDF -> 图片
        print(f"\n[INFO] 步骤 1/2: 将 PDF 转换为图片...")
        images = pdf_to_images(pdf_path, dpi=dpi)
        
        if not images:
            print("[ERROR] 未能从 PDF 中提取图片")
            return False
        
        # 第二步：图片 -> PPT
        print(f"\n[INFO] 步骤 2/2: 将图片插入 PowerPoint...")
        success = create_ppt_from_images(images, ppt_path)
        
        if not success:
            return False
        
        # 验证输出文件
        if not os.path.exists(ppt_path):
            print(f"[ERROR] PPTX 文件未生成: {ppt_path}")
            return False
        
        file_size = os.path.getsize(ppt_path)
        if file_size > 0:
            elapsed = time.time() - start_time
            print(f"\n[SUCCESS] ========== 转换完成 ==========")
            print(f"[INFO] 输出文件: {ppt_path}")
            print(f"[INFO] 文件大小: {file_size / (1024*1024):.2f} MB")
            print(f"[INFO] 总页数: {len(images)} 页")
            print(f"[INFO] 总耗时: {elapsed:.1f} 秒")
            print(f"[INFO] 平均速度: {len(images)/elapsed:.1f} 页/秒")
            return True
        else:
            print("[ERROR] PPTX 文件为空")
            return False
            
    except Exception as e:
        print(f"[ERROR] 转换失败: {str(e)}")
        traceback.print_exc()
        return False


def main():
    parser = argparse.ArgumentParser(
        description="PDF 转 PowerPoint - 图像级转换（每页渲染为高质量图片）"
    )
    parser.add_argument("-i", "--input", required=True, help="输入 PDF 文件路径")
    parser.add_argument("-o", "--output", required=True, help="输出 PowerPoint 文件路径")
    parser.add_argument("--dpi", type=int, default=200, 
                       help="图像分辨率 DPI (默认: 200, 推荐: 150-300)")

    args = parser.parse_args()

    # 验证输入文件
    if not os.path.exists(args.input):
        print(f"[ERROR] 输入文件不存在: {args.input}", file=sys.stderr)
        sys.exit(1)

    if not args.input.lower().endswith(".pdf"):
        print("[ERROR] 输入文件必须是 PDF 格式", file=sys.stderr)
        sys.exit(1)

    print(f"[INFO] 输入文件: {args.input}")
    print(f"[INFO] 输出文件: {args.output}")
    print(f"[INFO] 图像质量: {args.dpi} DPI")

    success = pdf_to_ppt(args.input, args.output, dpi=args.dpi)

    if success and os.path.exists(args.output):
        output_size = os.path.getsize(args.output)
        if output_size > 0:
            print(f"\n[SUCCESS] 转换成功!")
            print(f"[INFO] 输出文件: {args.output}")
            print(f"[INFO] 文件大小: {output_size / (1024*1024):.2f} MB")
            sys.exit(0)
        else:
            print("[ERROR] 输出文件为空", file=sys.stderr)
            sys.exit(1)
    else:
        print("[ERROR] 转换失败", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
